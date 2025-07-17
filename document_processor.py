import os
from typing import List, Dict, Any
import PyPDF2
from pptx import Presentation
import openpyxl
from io import BytesIO
from config import SUPPORTED_FILE_TYPES

class DocumentProcessor:
    def __init__(self):
        # Simple initialization
        pass
    
    def process_pdf(self, file_content: bytes) -> List[Dict[str, Any]]:
        """Process PDF files into chunks with one chunk per page"""
        chunks = []
        # Convert bytes to BytesIO object
        pdf_file = BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        for page_num, page in enumerate(pdf_reader.pages, 1):
            # Extract text for the entire page
            text = page.extract_text()
            
            # Create a chunk for each page
            chunks.append({
                "text": text,
                "metadata": {
                    "page": page_num,
                    "source": "pdf"
                }
            })
        return chunks

    def process_pptx(self, file_content: bytes) -> List[Dict[str, Any]]:
        """Process PowerPoint files into chunks with one chunk per slide"""
        chunks = []
        # Convert bytes to BytesIO object
        pptx_file = BytesIO(file_content)
        presentation = Presentation(pptx_file)
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            # Join all text from the slide
            full_text = ' '.join(slide_text)
            
            if full_text.strip():
                # Create a chunk for each slide
                chunks.append({
                    "text": full_text,
                    "metadata": {
                        "page": slide_num,
                        "source": "pptx"
                    }
                })
        return chunks

    def process_docx(self, file_content: bytes) -> List[Dict[str, Any]]:
        """Process Word documents into chunks with one chunk per page"""
        from docx import Document
        import docx
        
        # Convert bytes to BytesIO object
        docx_file = BytesIO(file_content)
        doc = Document(docx_file)
        chunks = []
        
        # We'll estimate pages based on page breaks and section breaks
        # Initialize variables to track pages
        current_page = 1
        current_page_text = []
        
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            
            # Check if this paragraph contains a page break
            has_page_break = False
            for run in paragraph.runs:
                if run.element.xpath('.//w:br[@w:type="page"]'):
                    has_page_break = True
                    break
            
            # Add text to current page
            if text:
                current_page_text.append(text)
            
            # If we found a page break or a section break, create a chunk for the current page
            if has_page_break or any(isinstance(child, docx.oxml.section.CT_SectPr) for child in paragraph._element.getchildren()):
                if current_page_text:
                    chunks.append({
                        "text": ' '.join(current_page_text),
                        "metadata": {
                            "page": current_page,
                            "source": "docx"
                        }
                    })
                    current_page += 1
                    current_page_text = []
        
        # Don't forget the last page if it has content
        if current_page_text:
            chunks.append({
                "text": ' '.join(current_page_text),
                "metadata": {
                    "page": current_page,
                    "source": "docx"
                }
            })
        
        return chunks

    def process_xlsx(self, file_content: bytes) -> List[Dict[str, Any]]:
        """Process Excel files into a single chunk for the entire workbook"""
        # Convert bytes to BytesIO object
        xlsx_file = BytesIO(file_content)
        workbook = openpyxl.load_workbook(xlsx_file)
        
        # Collect text from all sheets
        all_sheet_data = []
        
        for sheet_num, sheet_name in enumerate(workbook.sheetnames, 1):
            worksheet = workbook[sheet_name]
            sheet_data = []
            
            # Process each row in the worksheet
            for row in worksheet.iter_rows(values_only=True):
                row_text = ' '.join(str(cell) for cell in row if cell is not None)
                sheet_data.append(row_text)
            
            # Join sheet data
            sheet_full_text = ' '.join(sheet_data)
            all_sheet_data.append(sheet_full_text)
        
        # Create a single chunk for the entire workbook
        full_workbook_text = ' '.join(all_sheet_data)
        
        return [{
            "text": full_workbook_text,
            "metadata": {
                "source": "xlsx",
                "sheet_name": "workbook"
            }
        }]

    def process_document(self, file_content: bytes, file_type: str) -> List[Dict[str, Any]]:
        """Process documents based on file type"""
        processors = {
            "pdf": self.process_pdf,
            "pptx": self.process_pptx,
            "docx": self.process_docx,
            "xlsx": self.process_xlsx
        }
        
        if file_type not in processors:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        return processors[file_type](file_content)





# import os
# from typing import List, Dict, Any
# import PyPDF2
# from pptx import Presentation
# import openpyxl
# from io import BytesIO
# from config import SUPPORTED_FILE_TYPES

# class DocumentProcessor:
#     def __init__(self):
#         # Simple initialization
#         pass
    
#     def process_pdf(self, file_content: bytes) -> List[Dict[str, Any]]:
#         """Process PDF files into chunks with one chunk per page"""
#         chunks = []
#         # Convert bytes to BytesIO object
#         pdf_file = BytesIO(file_content)
#         pdf_reader = PyPDF2.PdfReader(pdf_file)
        
#         for page_num, page in enumerate(pdf_reader.pages, 1):
#             # Extract text for the entire page
#             text = page.extract_text()
            
#             # Create a chunk for each page
#             chunks.append({
#                 "text": text,
#                 "metadata": {
#                     "page": page_num,
#                     "source": "pdf"
#                 }
#             })
#         return chunks

#     def process_pptx(self, file_content: bytes) -> List[Dict[str, Any]]:
#         """Process PowerPoint files into chunks with one chunk per slide"""
#         chunks = []
#         # Convert bytes to BytesIO object
#         pptx_file = BytesIO(file_content)
#         presentation = Presentation(pptx_file)
        
#         for slide_num, slide in enumerate(presentation.slides, 1):
#             slide_text = []
#             for shape in slide.shapes:
#                 if hasattr(shape, "text"):
#                     slide_text.append(shape.text)
            
#             # Join all text from the slide
#             full_text = ' '.join(slide_text)
            
#             if full_text.strip():
#                 # Create a chunk for each slide
#                 chunks.append({
#                     "text": full_text,
#                     "metadata": {
#                         "page": slide_num,
#                         "source": "pptx"
#                     }
#                 })
#         return chunks

#     def process_docx(self, file_content: bytes) -> List[Dict[str, Any]]:
#         """Process Word documents into chunks with one chunk per page"""
#         from docx import Document
#         import docx
        
#         # Convert bytes to BytesIO object
#         docx_file = BytesIO(file_content)
#         doc = Document(docx_file)
#         chunks = []
        
#         # We'll estimate pages based on page breaks and section breaks
#         # Initialize variables to track pages
#         current_page = 1
#         current_page_text = []
        
#         for paragraph in doc.paragraphs:
#             text = paragraph.text.strip()
            
#             # Check if this paragraph contains a page break
#             has_page_break = False
#             for run in paragraph.runs:
#                 if run.element.xpath('.//w:br[@w:type="page"]'):
#                     has_page_break = True
#                     break
            
#             # Add text to current page
#             if text:
#                 current_page_text.append(text)
            
#             # If we found a page break or a section break, create a chunk for the current page
#             if has_page_break or any(isinstance(child, docx.oxml.section.CT_SectPr) for child in paragraph._element.getchildren()):
#                 if current_page_text:
#                     chunks.append({
#                         "text": ' '.join(current_page_text),
#                         "metadata": {
#                             "page": current_page,
#                             "source": "docx"
#                         }
#                     })
#                     current_page += 1
#                     current_page_text = []
        
#         # Don't forget the last page if it has content
#         if current_page_text:
#             chunks.append({
#                 "text": ' '.join(current_page_text),
#                 "metadata": {
#                     "page": current_page,
#                     "source": "docx"
#                 }
#             })
        
#         return chunks

#     def process_xlsx(self, file_content: bytes) -> List[Dict[str, Any]]:
#         """Process Excel files with one chunk per sheet"""
#         # Convert bytes to BytesIO object
#         xlsx_file = BytesIO(file_content)
#         workbook = openpyxl.load_workbook(xlsx_file)
        
#         chunks = []
        
#         # Process each sheet separately
#         for sheet_num, sheet_name in enumerate(workbook.sheetnames, 1):
#             worksheet = workbook[sheet_name]
#             sheet_data = []
            
#             # Process each row in the worksheet
#             for row in worksheet.iter_rows(values_only=True):
#                 row_text = ' '.join(str(cell) for cell in row if cell is not None)
#                 if row_text.strip():  # Only add non-empty rows
#                     sheet_data.append(row_text)
            
#             # Join sheet data
#             sheet_full_text = ' '.join(sheet_data)
            
#             # Create a chunk for this sheet
#             if sheet_full_text.strip():  # Only add non-empty sheets
#                 chunks.append({
#                     "text": sheet_full_text,
#                     "metadata": {
#                         "source": "xlsx",
#                         "sheet_name": sheet_name,
#                         "page": sheet_num  # Use sheet number as page for consistency
#                     }
#                 })
        
#         # If no valid chunks were created, create a default empty one
#         if not chunks:
#             chunks.append({
#                 "text": "No data found in Excel file.",
#                 "metadata": {
#                     "source": "xlsx",
#                     "sheet_name": "unknown",
#                     "page": 1
#                 }
#             })
            
#         return chunks

#     def process_document(self, file_content: bytes, file_type: str) -> List[Dict[str, Any]]:
#         """Process documents based on file type"""
#         processors = {
#             "pdf": self.process_pdf,
#             "pptx": self.process_pptx,
#             "docx": self.process_docx,
#             "xlsx": self.process_xlsx
#         }
        
#         if file_type not in processors:
#             raise ValueError(f"Unsupported file type: {file_type}")
        
#         return processors[file_type](file_content)
